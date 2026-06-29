import os
import io
import re
from flask import Flask, request, send_file, jsonify
from pptx import Presentation

app = Flask(__name__)
TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "template.pptx")


def fill_text_frame(text_frame, data):
    """Replace {{tag}} placeholders even when split across multiple runs."""
    for para in text_frame.paragraphs:
        if not para.runs:
            continue
        full_text = "".join(run.text for run in para.runs)
        if "{{" not in full_text:
            continue

        def replace_tag(match):
            key = match.group(1)
            return str(data.get(key, ""))

        new_text = re.sub(r"\{\{(\w+)\}\}", replace_tag, full_text)
        if new_text != full_text:
            para.runs[0].text = new_text
            for run in para.runs[1:]:
                run.text = ""


def fill_table(table, data):
    for row in table.rows:
        for cell in row.cells:
            fill_text_frame(cell.text_frame, data)


def remove_fully_empty_rows(prs):
    """
    After filling, scan every table on every slide. For any row (other than
    the first row, assumed to be a header) where ALL cells are blank once
    trimmed, delete that row. Works generically across any table — spec
    tables, ordering tables, lane-wavelength tables, etc.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            rows = list(table.rows)
            if len(rows) < 2:
                continue  # nothing to trim if there's only a header row
            for tr in rows[1:][::-1]:
                cell_texts = [c.text.strip() for c in tr.cells]
                if all(t == "" for t in cell_texts):
                    tr._tr.getparent().remove(tr._tr)


def remove_empty_tables(prs):
    """
    After row-removal, delete any table shape left with only its header row
    (or no rows at all) — meaning every data row in it was empty and got
    stripped out. This handles cases like an L4-L7 lane table that's
    entirely unused for a 4-lane product.
    """
    for slide in prs.slides:
        shapes_to_remove = []
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            if len(table.rows) <= 1:
                shapes_to_remove.append(shape)
        for shape in shapes_to_remove:
            shape._element.getparent().remove(shape._element)


def fill_template(data):
    prs = Presentation(TEMPLATE_PATH)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                fill_text_frame(shape.text_frame, data)
            if shape.has_table:
                fill_table(shape.table, data)
    remove_fully_empty_rows(prs)
    remove_empty_tables(prs)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


@app.route("/generate", methods=["POST"])
def generate():
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"error": "No JSON body received"}), 400
    try:
        file_buf = fill_template(data)
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    filename = f"{data.get('part_number_1', 'datasheet')}.pptx"
    return send_file(
        file_buf,
        as_attachment=True,
        download_name=filename,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok"})


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 5000)))